"""
Real engine tests — every test here runs an actual file through the
actual library/CLI tool, not a mock. Tests that depend on a system binary
(LibreOffice, ffmpeg, poppler, tesseract, the Playwright Chromium build)
skip cleanly when that binary isn't present, rather than failing, so this
suite still runs on a machine that only has the Python deps installed —
but `infrastructure/docker/Dockerfile.api` and `.github/workflows/test.yml`
both install the full toolchain, so in CI and in the real image these
run for real.
"""
import io
import os
import shutil
import sys
import zipfile

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from services.engines import (  # noqa: E402
    DocumentConvertEngine,
    MediaConvertEngine,
    PdfGenerateEngine,
    PdfManipulateEngine,
)

needs = lambda binary: pytest.mark.skipif(  # noqa: E731
    shutil.which(binary) is None, reason=f"{binary} not installed"
)


def _rembg_model_host_reachable() -> bool:
    """image_bg_remove's model (u2netp.onnx) downloads from rembg's GitHub
    release on first use, then rembg caches it on disk (~/.rembg/models)
    for every run after that. Mirrors test_job_queue.py's needs_redis
    pattern — skip cleanly with a documented reason rather than fail when
    that first download can't happen, instead of pretending network
    access is guaranteed everywhere this suite runs."""
    import urllib.error
    import urllib.request

    url = "https://github.com/danielgatis/rembg/releases/download/v0.0.0/u2netp.onnx"
    try:
        urllib.request.urlopen(urllib.request.Request(url, method="HEAD"), timeout=5)
        return True
    except urllib.error.URLError:
        return False


needs_rembg_model = pytest.mark.skipif(
    not _rembg_model_host_reachable(),
    reason="rembg's model host (github.com releases) not reachable — see docs/engines.md",
)


# -- PdfManipulateEngine (pure Python — pikepdf/pypdf/reportlab, no subprocess) --

class TestPdfManipulateEngine:
    def setup_method(self):
        self.engine = PdfManipulateEngine()

    def _run(self, tool_name, data, **options):
        return self.engine.process(io.BytesIO(data), {"tool_name": tool_name, **options})

    def test_merge(self, sample_pdf_bytes):
        result = self._run("pdf_merge", sample_pdf_bytes, extra_files=[sample_pdf_bytes])
        assert result.ok, result.error
        assert result.meta["page_count"] == 4  # 2 + 2 pages

    def test_merge_without_extra_files_fails_cleanly(self, sample_pdf_bytes):
        result = self._run("pdf_merge", sample_pdf_bytes)
        assert result.ok is False
        assert "extra_files" in result.error
        assert result.refundable is False  # missing required file — the user's mistake, not ours

    def test_split_default_one_pdf_per_page(self, sample_pdf_bytes):
        result = self._run("pdf_split", sample_pdf_bytes)
        assert result.ok, result.error
        with zipfile.ZipFile(io.BytesIO(result.output_bytes)) as zf:
            assert len(zf.namelist()) == 2

    def test_extract_pages(self, sample_pdf_bytes):
        result = self._run("pdf_extract_pages", sample_pdf_bytes, pages="2")
        assert result.ok, result.error
        assert result.meta["extracted"] == 1

    def test_extract_pages_requires_pages_option(self, sample_pdf_bytes):
        result = self._run("pdf_extract_pages", sample_pdf_bytes)
        assert result.ok is False
        assert result.refundable is False  # missing required option — the user's mistake, not ours

    def test_watermark_produces_valid_pdf(self, sample_pdf_bytes):
        from pypdf import PdfReader

        result = self._run("pdf_watermark", sample_pdf_bytes, text="CONFIDENTIAL")
        assert result.ok, result.error
        reader = PdfReader(io.BytesIO(result.output_bytes))
        assert len(reader.pages) == 2

    def test_rotate(self, sample_pdf_bytes):
        from pypdf import PdfReader

        result = self._run("pdf_rotate", sample_pdf_bytes, angle=90)
        assert result.ok, result.error
        reader = PdfReader(io.BytesIO(result.output_bytes))
        assert reader.pages[0].rotation == 90

    def test_rotate_rejects_non_multiple_of_90(self, sample_pdf_bytes):
        result = self._run("pdf_rotate", sample_pdf_bytes, angle=45)
        assert result.ok is False
        assert result.refundable is False  # invalid option value — the user's mistake, not ours

    def test_protect_then_unlock_round_trip(self, sample_pdf_bytes):
        from pypdf import PdfReader

        protected = self._run("pdf_protect", sample_pdf_bytes, password="secret123")
        assert protected.ok, protected.error
        reader = PdfReader(io.BytesIO(protected.output_bytes))
        assert reader.is_encrypted

        unlocked = self._run("pdf_unlock", protected.output_bytes, password="secret123")
        assert unlocked.ok, unlocked.error
        reader2 = PdfReader(io.BytesIO(unlocked.output_bytes))
        assert not reader2.is_encrypted

    def test_unlock_wrong_password_fails(self, sample_pdf_bytes):
        protected = self._run("pdf_protect", sample_pdf_bytes, password="secret123")
        result = self._run("pdf_unlock", protected.output_bytes, password="wrong")
        assert result.ok is False
        # A wrong password is the user's mistake, not TweakHub's — no refund.
        assert result.refundable is False

    def test_compress_produces_valid_pdf(self, sample_pdf_bytes):
        result = self._run("pdf_compress", sample_pdf_bytes)
        assert result.ok, result.error
        assert "bytes_after" in result.meta

    def test_repair_on_valid_pdf(self, sample_pdf_bytes):
        result = self._run("pdf_repair", sample_pdf_bytes)
        assert result.ok, result.error

    def test_organize_reorders_pages(self, sample_pdf_bytes):
        from pypdf import PdfReader

        original = PdfReader(io.BytesIO(sample_pdf_bytes))
        page2_text = original.pages[1].extract_text()

        result = self._run("pdf_organize", sample_pdf_bytes, order=[2, 1])
        assert result.ok, result.error
        reordered = PdfReader(io.BytesIO(result.output_bytes))
        assert reordered.pages[0].extract_text() == page2_text

    def test_sign_redact_pdfa_are_documented_stubs(self, sample_pdf_bytes):
        for tool in ("pdf_sign", "pdf_redact", "pdf_to_pdfa"):
            result = self._run(tool, sample_pdf_bytes)
            assert result.ok is False
            assert "Not implemented" in result.error
            # TweakHub hasn't built this yet — that's our doing, not the
            # user's mistake, so still refundable.
            assert result.refundable is True

    def test_corrupted_pdf_input_is_not_refundable(self):
        # pypdf can't parse this at all (PdfReadError/PyPdfError) — a
        # malformed file the user uploaded, not a TweakHub bug.
        result = self._run("pdf_rotate", b"not a real pdf file at all", angle=90)
        assert result.ok is False
        assert result.refundable is False


# -- PdfGenerateEngine (pure Python — reportlab) --

class TestPdfGenerateEngine:
    def setup_method(self):
        self.engine = PdfGenerateEngine()

    def _run(self, tool_name, payload):
        import json

        return self.engine.process(io.BytesIO(json.dumps(payload).encode()), {"tool_name": tool_name})

    def test_invoice_generator(self):
        from pypdf import PdfReader

        result = self._run("invoice_generator", {
            "invoiceNumber": "INV-001", "issueDate": "2026-08-31", "dueDate": "2026-09-15",
            "currency": "USD",
            "lineItems": [{"description": "Widget", "quantity": 2, "unitPrice": 9.99}],
        })
        assert result.ok, result.error
        assert result.meta["total"] == pytest.approx(19.98)
        reader = PdfReader(io.BytesIO(result.output_bytes))
        assert len(reader.pages) == 1

    def test_certificate_generator(self):
        result = self._run("certificate_generator", {"recipientName": "Jane Doe", "date": "2026-08-31"})
        assert result.ok, result.error

    def test_report_generator_paginates_long_content(self):
        sections = [{"heading": f"Section {i}", "body": "word " * 400} for i in range(10)]
        result = self._run("report_generator", {"title": "Q3 Report", "sections": sections})
        assert result.ok, result.error
        assert result.meta["sections"] == 10

    def test_invalid_json_fails_cleanly(self):
        result = self.engine.process(io.BytesIO(b"not json"), {"tool_name": "invoice_generator"})
        assert result.ok is False
        assert result.refundable is False  # malformed input payload — the user's mistake, not ours


# -- MediaConvertEngine: pure-Python paths (Pillow / pypdf / openpyxl) --

class TestMediaConvertEngineImages:
    def setup_method(self):
        self.engine = MediaConvertEngine()

    def _run(self, tool_name, data, **options):
        return self.engine.process(io.BytesIO(data), {"tool_name": tool_name, **options})

    def test_image_convert_png_to_jpeg(self, sample_png_bytes):
        result = self._run("image_convert", sample_png_bytes, target_format="jpg")
        assert result.ok, result.error
        assert result.content_type == "image/jpeg"

    def test_image_convert_requires_target_format(self, sample_png_bytes):
        result = self._run("image_convert", sample_png_bytes)
        assert result.ok is False

    def test_image_resize(self, sample_png_bytes):
        from PIL import Image

        result = self._run("image_resize", sample_png_bytes, width=200)
        assert result.ok, result.error
        img = Image.open(io.BytesIO(result.output_bytes))
        assert img.size[0] == 200

    def test_image_crop_default_centered(self, sample_png_bytes):
        result = self._run("image_crop", sample_png_bytes)
        assert result.ok, result.error

    def test_image_watermark(self, sample_png_bytes):
        result = self._run("image_watermark", sample_png_bytes, text="SAMPLE")
        assert result.ok, result.error

    def test_image_rotate(self, sample_png_bytes):
        from PIL import Image

        orig = Image.open(io.BytesIO(sample_png_bytes))
        result = self._run("image_rotate", sample_png_bytes, angle=90)
        assert result.ok, result.error
        rotated = Image.open(io.BytesIO(result.output_bytes))
        assert rotated.size == (orig.size[1], orig.size[0])

    def test_svg_to_png(self):
        svg = b'<svg xmlns="http://www.w3.org/2000/svg" width="50" height="50"><circle cx="25" cy="25" r="20" fill="red"/></svg>'
        result = self._run("svg_to_png", svg)
        assert result.ok, result.error
        assert result.content_type == "image/png"

    def test_png_to_pdf(self, sample_png_bytes):
        result = self._run("png_to_pdf", sample_png_bytes)
        assert result.ok, result.error
        assert result.content_type == "application/pdf"
        assert result.output_bytes[:4] == b"%PDF"

    @pytest.mark.parametrize("fmt", ["webp", "gif", "bmp", "tiff"])
    def test_image_to_pdf_named_pairs(self, fmt):
        result = self._run(f"{fmt}_to_pdf", self._encode(fmt))
        assert result.ok, result.error
        assert result.content_type == "application/pdf"
        assert result.output_bytes[:4] == b"%PDF"

    # -- Specific named format pairs (tools_catalog.py) — each reuses
    # _image_convert with engine_op supplying target_format (see
    # test_tool_router.py for that auto-population itself); these tests
    # confirm the actual conversion works for real, per pair, rather than
    # trusting "image_convert works so this probably does too." --

    def _encode(self, fmt: str) -> bytes:
        from PIL import Image

        img = Image.new("RGB", (40, 30), (10, 200, 90))
        if fmt.lower() in ("jpg", "jpeg", "bmp"):
            img = img.convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format={"jpg": "JPEG"}.get(fmt.lower(), fmt.upper()))
        return buf.getvalue()

    @pytest.mark.parametrize(
        "tool_name,src_fmt,expected_content_type",
        [
            ("png_to_jpg", "png", "image/jpeg"),
            ("jpg_to_png", "jpg", "image/png"),
            ("png_to_webp", "png", "image/webp"),
            ("webp_to_png", "webp", "image/png"),
            ("bmp_to_png", "bmp", "image/png"),
            ("png_to_bmp", "png", "image/bmp"),
            ("tiff_to_png", "tiff", "image/png"),
            ("png_to_tiff", "png", "image/tiff"),
            ("gif_to_png", "gif", "image/png"),
            # Remaining pairs among the same six formats (tools_catalog.py) —
            # filling out the ordered-pair grid, same handler, same discipline.
            ("jpg_to_webp", "jpg", "image/webp"),
            ("jpg_to_gif", "jpg", "image/gif"),
            ("jpg_to_bmp", "jpg", "image/bmp"),
            ("jpg_to_tiff", "jpg", "image/tiff"),
            ("png_to_gif", "png", "image/gif"),
            ("webp_to_jpg", "webp", "image/jpeg"),
            ("webp_to_gif", "webp", "image/gif"),
            ("webp_to_bmp", "webp", "image/bmp"),
            ("webp_to_tiff", "webp", "image/tiff"),
            ("gif_to_jpg", "gif", "image/jpeg"),
            ("gif_to_webp", "gif", "image/webp"),
            ("gif_to_bmp", "gif", "image/bmp"),
            ("gif_to_tiff", "gif", "image/tiff"),
            ("bmp_to_jpg", "bmp", "image/jpeg"),
            ("bmp_to_webp", "bmp", "image/webp"),
            ("bmp_to_gif", "bmp", "image/gif"),
            ("bmp_to_tiff", "bmp", "image/tiff"),
            ("tiff_to_jpg", "tiff", "image/jpeg"),
            ("tiff_to_webp", "tiff", "image/webp"),
            ("tiff_to_gif", "tiff", "image/gif"),
            ("tiff_to_bmp", "tiff", "image/bmp"),
            # New formats — ICO (favicons) and AVIF (modern web images),
            # both confirmed working in this Pillow build (see
            # IMAGE_FORMAT_MIME's comment in media_convert.py).
            ("png_to_ico", "png", "image/x-icon"),
            ("ico_to_png", "ico", "image/png"),
            ("jpg_to_ico", "jpg", "image/x-icon"),
            ("png_to_avif", "png", "image/avif"),
            ("avif_to_png", "avif", "image/png"),
            ("jpg_to_avif", "jpg", "image/avif"),
            ("avif_to_jpg", "avif", "image/jpeg"),
            ("webp_to_avif", "webp", "image/avif"),
            ("avif_to_webp", "avif", "image/webp"),
        ],
    )
    def test_named_image_format_pair(self, tool_name, src_fmt, expected_content_type):
        from PIL import Image

        target_format = tool_name.rsplit("_to_", 1)[1]
        src_bytes = self._encode(src_fmt)

        result = self._run(tool_name, src_bytes, target_format=target_format)
        assert result.ok, result.error
        assert result.content_type == expected_content_type

        # Not just "didn't error" — the output must actually decode as a
        # real image in the requested format.
        out = Image.open(io.BytesIO(result.output_bytes))
        out.load()
        assert out.format.lower() in (target_format.lower(), {"jpg": "jpeg"}.get(target_format.lower()))

    @needs_rembg_model
    def test_image_bg_remove_makes_the_background_transparent(self):
        """A purpose-built subject-on-plain-background image (not
        sample_png_bytes, which is solid-color text with no real salient
        object) so this is a meaningful segmentation check, not just
        "didn't crash": the corner (background) must end up transparent
        and the shape's center (subject) must stay opaque."""
        from PIL import Image, ImageDraw

        img = Image.new("RGB", (200, 200), (255, 255, 255))
        draw = ImageDraw.Draw(img)
        draw.ellipse((40, 40, 160, 160), fill=(220, 30, 30))
        buf = io.BytesIO()
        img.save(buf, format="PNG")

        result = self._run("image_bg_remove", buf.getvalue())
        assert result.ok, result.error
        assert result.content_type == "image/png"

        out = Image.open(io.BytesIO(result.output_bytes)).convert("RGBA")
        assert out.getpixel((5, 5))[3] < 50  # corner (background) — mostly transparent
        assert out.getpixel((100, 100))[3] > 200  # center (subject) — mostly opaque


class TestMediaConvertEnginePdfText:
    def setup_method(self):
        self.engine = MediaConvertEngine()

    def test_pdf_to_text_extracts_real_text(self, sample_pdf_bytes):
        result = self.engine.process(io.BytesIO(sample_pdf_bytes), {"tool_name": "pdf_to_text"})
        assert result.ok, result.error
        assert b"TweakHub test document" in result.output_bytes

    def test_text_to_pdf(self):
        result = self.engine.process(io.BytesIO(b"Hello TweakHub"), {"tool_name": "text_to_pdf"})
        assert result.ok, result.error
        assert result.content_type == "application/pdf"

    def test_pdf_to_tiff(self, sample_pdf_bytes):
        # pdftoppm natively supports -tiff alongside -jpeg/-png (verified
        # with `pdftoppm -h` before wiring this up) — same _pdf_to_image
        # handler, one more entry in its format lookup.
        from PIL import Image

        result = self.engine.process(
            io.BytesIO(sample_pdf_bytes), {"tool_name": "pdf_to_tiff", "target_format": "tiff"}
        )
        assert result.ok, result.error
        # sample_pdf_bytes is 2 pages, so multi-page output is a zip —
        # same convention as pdf_to_jpg/pdf_to_png on a multi-page source.
        assert result.content_type == "application/zip"
        assert result.meta["pages"] == 2
        import zipfile

        with zipfile.ZipFile(io.BytesIO(result.output_bytes)) as z:
            names = z.namelist()
            assert len(names) == 2
            img = Image.open(io.BytesIO(z.read(names[0])))
            img.load()
            assert img.format == "TIFF"

    def test_pdf_compare_produces_diff(self, sample_pdf_bytes):
        from reportlab.pdfgen import canvas

        buf = io.BytesIO()
        c = canvas.Canvas(buf)
        c.drawString(100, 750, "A totally different document")
        c.save()
        different_pdf = buf.getvalue()

        result = self.engine.process(
            io.BytesIO(sample_pdf_bytes), {"tool_name": "pdf_compare", "extra_files": [different_pdf]}
        )
        assert result.ok, result.error
        assert result.meta["changed_lines"] > 0


class TestMediaConvertEngineSpreadsheets:
    def setup_method(self):
        self.engine = MediaConvertEngine()

    def test_csv_to_xlsx_round_trips_through_xlsx_to_csv(self):
        csv_bytes = b"name,value\nWidget,42\nGadget,7\n"
        xlsx_result = self.engine.process(io.BytesIO(csv_bytes), {"tool_name": "csv_to_xlsx"})
        assert xlsx_result.ok, xlsx_result.error
        assert xlsx_result.meta["rows"] == 3

        csv_result = self.engine.process(io.BytesIO(xlsx_result.output_bytes), {"tool_name": "xlsx_to_csv"})
        assert csv_result.ok, csv_result.error
        assert b"Widget" in csv_result.output_bytes
        assert b"42" in csv_result.output_bytes


# -- MediaConvertEngine: ffmpeg-backed (skips if ffmpeg isn't installed) --

@needs("ffmpeg")
class TestMediaConvertEngineAudioVideo:
    def setup_method(self):
        self.engine = MediaConvertEngine()

    @pytest.fixture()
    def tone_wav(self):
        from services.engines._util import run, scratch_dir

        with scratch_dir() as d:
            out = d / "tone.wav"
            run(["ffmpeg", "-y", "-f", "lavfi", "-i", "sine=frequency=440:duration=1", "-ac", "1", str(out)])
            yield out.read_bytes()

    @pytest.fixture()
    def tiny_video(self):
        from services.engines._util import run, scratch_dir

        with scratch_dir() as d:
            out = d / "clip.mp4"
            run([
                "ffmpeg", "-y", "-f", "lavfi", "-i", "testsrc=duration=1:size=160x120:rate=10",
                "-f", "lavfi", "-i", "sine=frequency=440:duration=1",
                "-c:v", "libx264", "-c:a", "aac", "-shortest", str(out),
            ])
            yield out.read_bytes()

    def test_audio_convert(self, tone_wav):
        result = self.engine.process(io.BytesIO(tone_wav), {"tool_name": "audio_convert", "target_format": "mp3"})
        assert result.ok, result.error
        assert result.content_type == "audio/mp3"

    def test_audio_trim(self, tone_wav):
        result = self.engine.process(io.BytesIO(tone_wav), {"tool_name": "audio_trim", "duration": "0.5"})
        assert result.ok, result.error

    def test_audio_normalize(self, tone_wav):
        result = self.engine.process(io.BytesIO(tone_wav), {"tool_name": "audio_normalize"})
        assert result.ok, result.error

    def test_video_compress(self, tiny_video):
        result = self.engine.process(io.BytesIO(tiny_video), {"tool_name": "video_compress"})
        assert result.ok, result.error

    def test_video_to_gif(self, tiny_video):
        result = self.engine.process(io.BytesIO(tiny_video), {"tool_name": "video_to_gif"})
        assert result.ok, result.error
        assert result.content_type == "image/gif"

    def test_video_extract_audio(self, tiny_video):
        result = self.engine.process(io.BytesIO(tiny_video), {"tool_name": "video_extract_audio"})
        assert result.ok, result.error

    def test_video_mute(self, tiny_video):
        result = self.engine.process(io.BytesIO(tiny_video), {"tool_name": "video_mute"})
        assert result.ok, result.error

    # -- Specific named video/audio format pairs (tools_catalog.py). Each
    # reuses _video_convert/_audio_convert with engine_op supplying
    # target_format. The fixture's actual container (tiny_video is real
    # mp4, tone_wav is real wav) doesn't have to match a pair's nominal
    # source format (e.g. "mkv_to_mp4") — ffmpeg identifies real content
    # by its header, not the scratch file's extension (verified
    # separately by hand; see docs/TODO.md's per-tool-timeout-tuning
    # entry for where that finding is used), and every one of these pairs
    # was independently confirmed against its actual nominal source
    # format via manual ffmpeg runs before being added to the catalog. --

    @pytest.mark.parametrize(
        "tool_name,expected_content_type",
        [
            ("mp4_to_webm", "video/webm"),
            ("webm_to_mp4", "video/mp4"),
            ("mp4_to_mkv", "video/mkv"),
            ("mkv_to_mp4", "video/mp4"),
            ("mp4_to_mov", "video/mov"),
            ("mov_to_mp4", "video/mp4"),
            # Remaining pairs among mp4/webm/mkv/mov, plus two more
            # containers (avi, flv) bidirectional with mp4.
            ("webm_to_mkv", "video/mkv"),
            ("mkv_to_webm", "video/webm"),
            ("webm_to_mov", "video/mov"),
            ("mov_to_webm", "video/webm"),
            ("mkv_to_mov", "video/mov"),
            ("mov_to_mkv", "video/mkv"),
            ("mp4_to_avi", "video/avi"),
            ("avi_to_mp4", "video/mp4"),
            ("mp4_to_flv", "video/flv"),
            ("flv_to_mp4", "video/mp4"),
            # New containers — plain ffmpeg -i in -> out, no extra codec
            # args needed (3gp was tried and dropped — needs explicit
            # -c:v/-c:a args the generic handler doesn't pass).
            ("mp4_to_wmv", "video/wmv"),
            ("wmv_to_mp4", "video/mp4"),
            ("mp4_to_ts", "video/ts"),
            ("ts_to_mp4", "video/mp4"),
            ("mp4_to_m4v", "video/m4v"),
            ("m4v_to_mp4", "video/mp4"),
        ],
    )
    def test_named_video_format_pair(self, tiny_video, tool_name, expected_content_type):
        target_format = tool_name.rsplit("_to_", 1)[1]
        result = self.engine.process(
            io.BytesIO(tiny_video), {"tool_name": tool_name, "target_format": target_format}
        )
        assert result.ok, result.error
        assert result.content_type == expected_content_type
        assert len(result.output_bytes) > 0

    # -- Named "video to mp3" extraction pairs — same _video_extract_audio
    # handler as video_extract_audio, no target_format needed (always mp3).
    @pytest.mark.parametrize(
        "tool_name",
        [
            "mp4_to_mp3", "mov_to_mp3", "webm_to_mp3", "mkv_to_mp3", "avi_to_mp3",
            # Remaining containers — completes the extraction grid.
            "flv_to_mp3", "wmv_to_mp3", "ts_to_mp3", "m4v_to_mp3",
        ],
    )
    def test_named_video_to_mp3_pair(self, tiny_video, tool_name):
        result = self.engine.process(io.BytesIO(tiny_video), {"tool_name": tool_name})
        assert result.ok, result.error
        assert result.content_type == "audio/mpeg"
        assert len(result.output_bytes) > 0

    @pytest.mark.parametrize(
        "tool_name,expected_content_type",
        [
            ("mp3_to_wav", "audio/wav"),
            ("wav_to_mp3", "audio/mp3"),
            ("wav_to_flac", "audio/flac"),
            ("flac_to_mp3", "audio/mp3"),
            ("wav_to_ogg", "audio/ogg"),
            ("ogg_to_mp3", "audio/mp3"),
            ("wav_to_m4a", "audio/m4a"),
            ("m4a_to_mp3", "audio/mp3"),
            # Remaining pairs among mp3/wav/flac/ogg/m4a.
            ("mp3_to_flac", "audio/flac"),
            ("flac_to_wav", "audio/wav"),
            ("mp3_to_ogg", "audio/ogg"),
            ("ogg_to_wav", "audio/wav"),
            ("mp3_to_m4a", "audio/m4a"),
            ("m4a_to_wav", "audio/wav"),
            ("flac_to_ogg", "audio/ogg"),
            ("ogg_to_flac", "audio/flac"),
            ("flac_to_m4a", "audio/m4a"),
            ("m4a_to_flac", "audio/flac"),
            ("ogg_to_m4a", "audio/m4a"),
            ("m4a_to_ogg", "audio/ogg"),
            # New formats.
            ("mp3_to_opus", "audio/opus"),
            ("opus_to_mp3", "audio/mp3"),
            ("mp3_to_aac", "audio/aac"),
            ("aac_to_mp3", "audio/mp3"),
            ("mp3_to_wma", "audio/wma"),
            ("wma_to_mp3", "audio/mp3"),
            ("mp3_to_aiff", "audio/aiff"),
            ("aiff_to_mp3", "audio/mp3"),
            # Connects opus/aac/wma/aiff to the wav hub too, not just mp3.
            ("wav_to_opus", "audio/opus"),
            ("opus_to_wav", "audio/wav"),
            ("wav_to_aac", "audio/aac"),
            ("aac_to_wav", "audio/wav"),
            ("wav_to_wma", "audio/wma"),
            ("wma_to_wav", "audio/wav"),
            ("wav_to_aiff", "audio/aiff"),
            ("aiff_to_wav", "audio/wav"),
        ],
    )
    def test_named_audio_format_pair(self, tone_wav, tool_name, expected_content_type):
        target_format = tool_name.rsplit("_to_", 1)[1]
        result = self.engine.process(
            io.BytesIO(tone_wav), {"tool_name": tool_name, "target_format": target_format}
        )
        assert result.ok, result.error
        assert result.content_type == expected_content_type
        assert len(result.output_bytes) > 0


# -- DocumentConvertEngine: LibreOffice / poppler / tesseract / Playwright --

@needs("soffice")
class TestDocumentConvertEngineLibreOffice:
    def setup_method(self):
        self.engine = DocumentConvertEngine()

    def test_pdf_to_word(self, sample_pdf_bytes):
        result = self.engine.process(io.BytesIO(sample_pdf_bytes), {"tool_name": "pdf_to_word"})
        assert result.ok, result.error
        assert result.output_bytes[:2] == b"PK"  # docx is a zip

    def test_word_to_pdf_round_trip(self, sample_pdf_bytes):
        docx_result = self.engine.process(io.BytesIO(sample_pdf_bytes), {"tool_name": "pdf_to_word"})
        assert docx_result.ok, docx_result.error

        pdf_result = self.engine.process(io.BytesIO(docx_result.output_bytes), {"tool_name": "word_to_pdf"})
        assert pdf_result.ok, pdf_result.error
        assert pdf_result.output_bytes[:4] == b"%PDF"

    def test_pdf_to_excel_and_pdf_to_ppt_are_documented_stubs(self, sample_pdf_bytes):
        for tool in ("pdf_to_excel", "pdf_to_ppt"):
            result = self.engine.process(io.BytesIO(sample_pdf_bytes), {"tool_name": tool})
            assert result.ok is False
            assert "Not implemented" in result.error

    def test_docx_to_txt_extracts_real_text(self, sample_pdf_bytes):
        docx_result = self.engine.process(io.BytesIO(sample_pdf_bytes), {"tool_name": "pdf_to_word"})
        assert docx_result.ok, docx_result.error

        txt_result = self.engine.process(io.BytesIO(docx_result.output_bytes), {"tool_name": "docx_to_txt"})
        assert txt_result.ok, txt_result.error
        assert txt_result.content_type == "text/plain"
        assert len(txt_result.output_bytes.strip()) > 0

    def test_odt_to_docx_round_trip(self, sample_pdf_bytes):
        docx_result = self.engine.process(io.BytesIO(sample_pdf_bytes), {"tool_name": "pdf_to_word"})
        assert docx_result.ok, docx_result.error
        odt_result = self.engine.process(io.BytesIO(docx_result.output_bytes), {"tool_name": "docx_to_odt"})
        assert odt_result.ok, odt_result.error

        result = self.engine.process(io.BytesIO(odt_result.output_bytes), {"tool_name": "odt_to_docx"})
        assert result.ok, result.error
        assert result.output_bytes[:2] == b"PK"  # docx is a zip
        assert result.content_type == (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

    def test_xlsx_to_ods_and_back_preserves_data(self):
        from openpyxl import Workbook, load_workbook

        wb = Workbook()
        ws = wb.active
        ws.append(["name", "score"])
        ws.append(["alice", 42])
        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        ods_result = self.engine.process(io.BytesIO(xlsx_bytes), {"tool_name": "xlsx_to_ods"})
        assert ods_result.ok, ods_result.error
        assert ods_result.content_type == "application/vnd.oasis.opendocument.spreadsheet"

        xlsx2_result = self.engine.process(io.BytesIO(ods_result.output_bytes), {"tool_name": "ods_to_xlsx"})
        assert xlsx2_result.ok, xlsx2_result.error

        wb2 = load_workbook(io.BytesIO(xlsx2_result.output_bytes))
        rows = list(wb2.active.iter_rows(values_only=True))
        assert rows == [("name", "score"), ("alice", 42)]

    def test_pptx_to_odp_round_trip_and_odp_to_pdf(self):
        # A real minimal .pptx (python-pptx, test-fixture only — see
        # requirements.txt) rather than trusting docx/xlsx pairs already
        # working means presentations do too.
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "TweakHub test deck"
        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        odp_result = self.engine.process(io.BytesIO(pptx_bytes), {"tool_name": "pptx_to_odp"})
        assert odp_result.ok, odp_result.error
        assert odp_result.content_type == "application/vnd.oasis.opendocument.presentation"

        pptx2_result = self.engine.process(io.BytesIO(odp_result.output_bytes), {"tool_name": "odp_to_pptx"})
        assert pptx2_result.ok, pptx2_result.error
        assert pptx2_result.output_bytes[:2] == b"PK"  # pptx is a zip
        assert pptx2_result.content_type == (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        pdf_result = self.engine.process(io.BytesIO(odp_result.output_bytes), {"tool_name": "odp_to_pdf"})
        assert pdf_result.ok, pdf_result.error
        assert pdf_result.output_bytes[:4] == b"%PDF"

    def test_txt_to_docx(self):
        result = self.engine.process(
            io.BytesIO(b"Hello TweakHub, this is a plain text file.\nSecond line."),
            {"tool_name": "txt_to_docx"},
        )
        assert result.ok, result.error
        assert result.output_bytes[:2] == b"PK"  # docx is a zip
        assert result.content_type == (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

    def test_html_to_docx_and_docx_to_html_round_trip(self):
        html = b"<html><body><h1>TweakHub test</h1><p>Hello <b>world</b>.</p></body></html>"
        docx_result = self.engine.process(io.BytesIO(html), {"tool_name": "html_to_docx"})
        assert docx_result.ok, docx_result.error
        assert docx_result.output_bytes[:2] == b"PK"  # docx is a zip

        import zipfile

        with zipfile.ZipFile(io.BytesIO(docx_result.output_bytes)) as z:
            document_xml = z.read("word/document.xml").decode("utf-8", errors="replace")
        assert "TweakHub test" in document_xml
        assert "Hello" in document_xml

        html_result = self.engine.process(io.BytesIO(docx_result.output_bytes), {"tool_name": "docx_to_html"})
        assert html_result.ok, html_result.error
        assert html_result.content_type == "text/html"
        assert b"<html" in html_result.output_bytes.lower()
        assert b"TweakHub" in html_result.output_bytes

    def test_odt_to_txt_extracts_real_text(self):
        # Seeded via html_to_docx rather than pdf_to_word deliberately —
        # pdf_to_word's output is absolutely-positioned draw:custom-shape
        # text frames (how LibreOffice's PDF import inherently works),
        # which for a *multi-page* source loses everything but page 1
        # when re-exported as plain text. That's a real, documented
        # limitation (see this module's docstring), not something this
        # test should paper over — html_to_docx produces normal flowing
        # paragraphs, which is what odt_to_txt is actually meant for.
        html = b"<html><body><h1>TweakHub test document</h1><p>First paragraph.</p><p>Second paragraph.</p></body></html>"
        docx_result = self.engine.process(io.BytesIO(html), {"tool_name": "html_to_docx"})
        assert docx_result.ok, docx_result.error
        odt_result = self.engine.process(io.BytesIO(docx_result.output_bytes), {"tool_name": "docx_to_odt"})
        assert odt_result.ok, odt_result.error

        txt_result = self.engine.process(io.BytesIO(odt_result.output_bytes), {"tool_name": "odt_to_txt"})
        assert txt_result.ok, txt_result.error
        assert txt_result.content_type == "text/plain"
        assert b"TweakHub test document" in txt_result.output_bytes
        assert b"First paragraph" in txt_result.output_bytes
        assert b"Second paragraph" in txt_result.output_bytes

    def test_odt_to_txt_loses_content_past_page_1_from_a_multipage_pdf_to_word_chain(self, sample_pdf_bytes):
        """Documents the limitation itself (see module docstring) as a
        real, checked behavior rather than an unstated gotcha — so a
        future change that fixes it (a LibreOffice upgrade, say) is
        caught by this test failing, and a regression is caught too."""
        docx_result = self.engine.process(io.BytesIO(sample_pdf_bytes), {"tool_name": "pdf_to_word"})
        assert docx_result.ok, docx_result.error
        odt_result = self.engine.process(io.BytesIO(docx_result.output_bytes), {"tool_name": "docx_to_odt"})
        assert odt_result.ok, odt_result.error

        txt_result = self.engine.process(io.BytesIO(odt_result.output_bytes), {"tool_name": "odt_to_txt"})
        assert txt_result.ok, txt_result.error
        # sample_pdf_bytes is 2 pages ("...page 1" / "page 2") — page 2's
        # text does not survive this chain today.
        assert b"page 2" not in txt_result.output_bytes

    def test_rtf_to_docx(self, sample_pdf_bytes):
        # A real .rtf seed file, generated by LibreOffice itself (same
        # bootstrap pattern as every other chained seed file in this
        # suite) rather than hand-authoring RTF markup.
        from services.engines._util import run, scratch_dir

        docx_result = self.engine.process(io.BytesIO(sample_pdf_bytes), {"tool_name": "pdf_to_word"})
        assert docx_result.ok, docx_result.error

        with scratch_dir() as d:
            src = d / "in.docx"
            src.write_bytes(docx_result.output_bytes)
            out_dir = d / "out"
            out_dir.mkdir()
            run(["soffice", "--headless", "--convert-to", "rtf", "--outdir", str(out_dir), str(src)], timeout=90)
            rtf_bytes = (out_dir / "in.rtf").read_bytes()

        result = self.engine.process(io.BytesIO(rtf_bytes), {"tool_name": "rtf_to_docx"})
        assert result.ok, result.error
        assert result.output_bytes[:2] == b"PK"  # docx is a zip
        assert result.content_type == (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

    def test_doc_to_pdf_and_doc_to_docx_preserve_content(self):
        # A real legacy .doc seed file, generated by LibreOffice itself
        # exporting a normal (non-pdf_to_word-derived) docx — same
        # bootstrap pattern as rtf_to_docx above. Content is checked all
        # the way through, not just "produced a file" — burned once
        # already this pass by pdf_to_word's custom-shape quirk (see
        # module docstring), so these round trips are verified for real.
        from services.engines._util import run, scratch_dir

        html = b"<html><body><h1>TweakHub test document</h1><p>First paragraph of content.</p></body></html>"
        docx_result = self.engine.process(io.BytesIO(html), {"tool_name": "html_to_docx"})
        assert docx_result.ok, docx_result.error

        with scratch_dir() as d:
            src = d / "in.docx"
            src.write_bytes(docx_result.output_bytes)
            out_dir = d / "out"
            out_dir.mkdir()
            run([
                "soffice", f"-env:UserInstallation=file://{d}/lo_profile",
                "--headless", "--convert-to", "doc", "--outdir", str(out_dir), str(src),
            ], timeout=90)
            doc_bytes = (out_dir / "in.doc").read_bytes()

        pdf_result = self.engine.process(io.BytesIO(doc_bytes), {"tool_name": "doc_to_pdf"})
        assert pdf_result.ok, pdf_result.error
        assert pdf_result.output_bytes[:4] == b"%PDF"
        from pypdf import PdfReader

        text = PdfReader(io.BytesIO(pdf_result.output_bytes)).pages[0].extract_text()
        assert "TweakHub test document" in text
        assert "First paragraph of content" in text

        docx2_result = self.engine.process(io.BytesIO(doc_bytes), {"tool_name": "doc_to_docx"})
        assert docx2_result.ok, docx2_result.error
        with zipfile.ZipFile(io.BytesIO(docx2_result.output_bytes)) as z:
            document_xml = z.read("word/document.xml").decode("utf-8", errors="replace")
        assert "TweakHub test document" in document_xml
        assert "First paragraph of content" in document_xml

    def test_xls_to_xlsx_and_xls_to_pdf_preserve_data(self):
        from openpyxl import Workbook, load_workbook

        from services.engines._util import run, scratch_dir

        wb = Workbook()
        ws = wb.active
        ws.append(["name", "score"])
        ws.append(["alice", 42])
        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        with scratch_dir() as d:
            src = d / "seed.xlsx"
            src.write_bytes(xlsx_bytes)
            out_dir = d / "out"
            out_dir.mkdir()
            run([
                "soffice", f"-env:UserInstallation=file://{d}/lo_profile",
                "--headless", "--convert-to", "xls", "--outdir", str(out_dir), str(src),
            ], timeout=90)
            xls_bytes = (out_dir / "seed.xls").read_bytes()

        xlsx_result = self.engine.process(io.BytesIO(xls_bytes), {"tool_name": "xls_to_xlsx"})
        assert xlsx_result.ok, xlsx_result.error
        wb2 = load_workbook(io.BytesIO(xlsx_result.output_bytes))
        rows = list(wb2.active.iter_rows(values_only=True))
        assert rows == [("name", "score"), ("alice", 42)]

        pdf_result = self.engine.process(io.BytesIO(xls_bytes), {"tool_name": "xls_to_pdf"})
        assert pdf_result.ok, pdf_result.error
        assert pdf_result.output_bytes[:4] == b"%PDF"
        from pypdf import PdfReader

        text = PdfReader(io.BytesIO(pdf_result.output_bytes)).pages[0].extract_text()
        assert "alice" in text
        assert "42" in text

    def test_csv_to_pdf_preserves_data(self):
        csv_bytes = b"name,value\nWidget,42\nGadget,7\n"
        result = self.engine.process(io.BytesIO(csv_bytes), {"tool_name": "csv_to_pdf"})
        assert result.ok, result.error
        assert result.output_bytes[:4] == b"%PDF"
        from pypdf import PdfReader

        text = PdfReader(io.BytesIO(result.output_bytes)).pages[0].extract_text()
        assert "Widget" in text
        assert "42" in text

    def test_ods_to_pdf_preserves_data(self):
        # ODS seed generated by LibreOffice itself exporting a normal
        # xlsx — same bootstrap pattern as xls_to_xlsx above.
        from openpyxl import Workbook

        from services.engines._util import run, scratch_dir

        wb = Workbook()
        ws = wb.active
        ws.append(["name", "score"])
        ws.append(["alice", 42])
        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        with scratch_dir() as d:
            src = d / "seed.xlsx"
            src.write_bytes(xlsx_bytes)
            out_dir = d / "out"
            out_dir.mkdir()
            run([
                "soffice", f"-env:UserInstallation=file://{d}/lo_profile",
                "--headless", "--convert-to", "ods", "--outdir", str(out_dir), str(src),
            ], timeout=90)
            ods_bytes = (out_dir / "seed.ods").read_bytes()

        result = self.engine.process(io.BytesIO(ods_bytes), {"tool_name": "ods_to_pdf"})
        assert result.ok, result.error
        assert result.output_bytes[:4] == b"%PDF"
        from pypdf import PdfReader

        text = PdfReader(io.BytesIO(result.output_bytes)).pages[0].extract_text()
        assert "alice" in text
        assert "42" in text

    def test_rtf_to_txt_and_rtf_to_odt_preserve_content(self):
        # RTF seed generated by LibreOffice itself exporting a normal
        # (non-pdf_to_word-derived) docx — same bootstrap pattern used
        # throughout this batch.
        from services.engines._util import run, scratch_dir

        html = b"<html><body><h1>TweakHub test document</h1><p>First paragraph of content.</p></body></html>"
        docx_result = self.engine.process(io.BytesIO(html), {"tool_name": "html_to_docx"})
        assert docx_result.ok, docx_result.error

        with scratch_dir() as d:
            src = d / "in.docx"
            src.write_bytes(docx_result.output_bytes)
            out_dir = d / "out"
            out_dir.mkdir()
            run([
                "soffice", f"-env:UserInstallation=file://{d}/lo_profile",
                "--headless", "--convert-to", "rtf", "--outdir", str(out_dir), str(src),
            ], timeout=90)
            rtf_bytes = (out_dir / "in.rtf").read_bytes()

        txt_result = self.engine.process(io.BytesIO(rtf_bytes), {"tool_name": "rtf_to_txt"})
        assert txt_result.ok, txt_result.error
        assert b"TweakHub test document" in txt_result.output_bytes
        assert b"First paragraph of content" in txt_result.output_bytes

        odt_result = self.engine.process(io.BytesIO(rtf_bytes), {"tool_name": "rtf_to_odt"})
        assert odt_result.ok, odt_result.error
        with zipfile.ZipFile(io.BytesIO(odt_result.output_bytes)) as z:
            content_xml = z.read("content.xml").decode("utf-8", errors="replace")
        assert "TweakHub test document" in content_xml
        assert "First paragraph of content" in content_xml

    def test_doc_to_txt_preserves_content(self):
        from services.engines._util import run, scratch_dir

        html = b"<html><body><h1>TweakHub test document</h1><p>First paragraph of content.</p></body></html>"
        docx_result = self.engine.process(io.BytesIO(html), {"tool_name": "html_to_docx"})
        assert docx_result.ok, docx_result.error

        with scratch_dir() as d:
            src = d / "in.docx"
            src.write_bytes(docx_result.output_bytes)
            out_dir = d / "out"
            out_dir.mkdir()
            run([
                "soffice", f"-env:UserInstallation=file://{d}/lo_profile",
                "--headless", "--convert-to", "doc", "--outdir", str(out_dir), str(src),
            ], timeout=90)
            doc_bytes = (out_dir / "in.doc").read_bytes()

        result = self.engine.process(io.BytesIO(doc_bytes), {"tool_name": "doc_to_txt"})
        assert result.ok, result.error
        assert b"TweakHub test document" in result.output_bytes
        assert b"First paragraph of content" in result.output_bytes

    def test_xls_to_csv_preserves_data(self):
        from openpyxl import Workbook

        from services.engines._util import run, scratch_dir

        wb = Workbook()
        ws = wb.active
        ws.append(["name", "score"])
        ws.append(["alice", 42])
        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        with scratch_dir() as d:
            src = d / "seed.xlsx"
            src.write_bytes(xlsx_bytes)
            out_dir = d / "out"
            out_dir.mkdir()
            run([
                "soffice", f"-env:UserInstallation=file://{d}/lo_profile",
                "--headless", "--convert-to", "xls", "--outdir", str(out_dir), str(src),
            ], timeout=90)
            xls_bytes = (out_dir / "seed.xls").read_bytes()

        result = self.engine.process(io.BytesIO(xls_bytes), {"tool_name": "xls_to_csv"})
        assert result.ok, result.error
        assert b"alice" in result.output_bytes
        assert b"42" in result.output_bytes
        assert result.content_type == "text/csv"

    def test_ods_to_csv_and_csv_to_ods_preserve_data(self):
        # Same CSV export filter already verified for xls_to_csv,
        # re-verified here against a real ODS seed rather than assumed to
        # carry over — plus the reverse import direction.
        from openpyxl import Workbook

        from services.engines._util import run, scratch_dir

        wb = Workbook()
        ws = wb.active
        ws.append(["name", "score"])
        ws.append(["alice", 42])
        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        with scratch_dir() as d:
            src = d / "seed.xlsx"
            src.write_bytes(xlsx_bytes)
            out_dir = d / "out"
            out_dir.mkdir()
            run([
                "soffice", f"-env:UserInstallation=file://{d}/lo_profile",
                "--headless", "--convert-to", "ods", "--outdir", str(out_dir), str(src),
            ], timeout=90)
            ods_bytes = (out_dir / "seed.ods").read_bytes()

        csv_result = self.engine.process(io.BytesIO(ods_bytes), {"tool_name": "ods_to_csv"})
        assert csv_result.ok, csv_result.error
        assert b"alice" in csv_result.output_bytes
        assert b"42" in csv_result.output_bytes
        assert csv_result.content_type == "text/csv"

        csv_bytes = b"name,score\nalice,42\n"
        ods_result = self.engine.process(io.BytesIO(csv_bytes), {"tool_name": "csv_to_ods"})
        assert ods_result.ok, ods_result.error
        with zipfile.ZipFile(io.BytesIO(ods_result.output_bytes)) as z:
            content_xml = z.read("content.xml").decode("utf-8", errors="replace")
        assert "alice" in content_xml
        assert "42" in content_xml


@needs("pdftohtml")
class TestDocumentConvertEnginePoppler:
    def test_pdf_to_html(self, sample_pdf_bytes):
        engine = DocumentConvertEngine()
        result = engine.process(io.BytesIO(sample_pdf_bytes), {"tool_name": "pdf_to_html"})
        assert result.ok, result.error
        assert b"<html" in result.output_bytes.lower()


class TestDocumentConvertEnginePdfToMarkdown:
    def test_pdf_to_markdown(self, sample_pdf_bytes):
        engine = DocumentConvertEngine()
        result = engine.process(io.BytesIO(sample_pdf_bytes), {"tool_name": "pdf_to_markdown"})
        assert result.ok, result.error
        assert "## Page 1" in result.output_bytes.decode()


@needs("tesseract")
@needs("pdftoppm")
class TestDocumentConvertEngineOcr:
    def test_ocr_extract_produces_searchable_pdf(self, sample_pdf_bytes):
        engine = DocumentConvertEngine()
        result = engine.process(io.BytesIO(sample_pdf_bytes), {"tool_name": "ocr_extract"})
        assert result.ok, result.error
        assert result.meta["pages_ocred"] == 2


def _chromium_available() -> bool:
    from pathlib import Path

    base = Path(os.environ.get("PLAYWRIGHT_BROWSERS_PATH", "/opt/pw-browsers"))
    return any(base.glob("chromium-*/chrome-linux/chrome"))


@pytest.mark.skipif(not _chromium_available(), reason="Playwright Chromium not installed")
class TestDocumentConvertEngineHtmlToPdf:
    def test_html_to_pdf(self):
        engine = DocumentConvertEngine()
        html = b"<html><body><h1>Hello TweakHub</h1></body></html>"
        result = engine.process(io.BytesIO(html), {"tool_name": "html_to_pdf"})
        assert result.ok, result.error
        assert result.output_bytes[:4] == b"%PDF"

    def test_markdown_to_pdf(self):
        engine = DocumentConvertEngine()
        result = engine.process(io.BytesIO(b"# Hello\n\nSome **bold** text."), {"tool_name": "markdown_to_pdf"})
        assert result.ok, result.error
        assert result.output_bytes[:4] == b"%PDF"
